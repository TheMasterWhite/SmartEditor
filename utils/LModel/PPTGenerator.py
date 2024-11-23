import logging
import os.path
from .Interface import LLMInterface
from ..SModel.ImageGenerator import ImageGenerator
import erniebot, copy
from utils.Config.FileProcess import FileProcess, OSSProcess, JsonOperator
from Config import *
from utils import Tools
from pptx import Presentation

fileSavePath = copy.deepcopy(GLOBAL_FileSavePath)


# PPT页面操作器类
class SlideOperator:

    # 传入形状替换文本
    @staticmethod
    def replace_text(Shape, Content):
        if not Shape.has_text_frame:  # 判断是否有文本框
            return
        tf = Shape.text_frame
        for paragraph in tf.paragraphs:
            is_first_run = True
            for run in paragraph.runs:
                if is_first_run:
                    run.text = Content
                    is_first_run = False
                else:
                    run.text = ""


    # 替换PPT配图
    @staticmethod
    def replace_image(Slide, ImagePath, ShapeName):
        for shape in Slide.shapes:
            if shape.name == ShapeName:
                shape.image = ImagePath


    # 根据形状名字替换文本
    @staticmethod
    def replace_shape_text(Slide, ShapeName, Content):
        # 遍历所有图形，并替换对应章节
        for shape in Slide.shapes:
            # 替换目录内容
            if shape.name == ShapeName:
                SlideOperator.replace_text(shape, Content)


    # 根据下标删除某一页的幻灯片
    @staticmethod
    def delete_slide(Prs, Index):
        rid = Prs.slides._sldIdLst[Index].rId
        Prs.part.drop_rel(rid)
        del Prs.slides._sldIdLst[Index]
        return Prs


    # 插入图片
    @staticmethod
    def insert_image(Slide, PageSummary):
        # 遍历所有图形，并替换图片
        imageCount = 0
        for shape in Slide.shapes:
            # 替换目录内容
            if shape.name == "图片":
                imageCount = 1
                break
            elif shape.name == "图片1":
                imageCount = 3
                break
        if imageCount == 0:
            return

        promptPath = os.path.join(GLOBAL_ResourcesSavePath, "PPT配图.txt")
        with open(promptPath, "r", encoding = "utf-8") as f:
            prompt = f.read()

        if imageCount == 1:
            prompt += PageSummary["标题"]
            for content in PageSummary["内容"]:
                prompt += content
            imagePrompt = LLMInterface.GetResponse_String(prompt)
            imageUrl = ImageGenerator.generate_image(Prompt = imagePrompt, Size = (540, 540))
            savePath = ImageGenerator.download_image(imageUrl)
            ImageGenerator.resize_image(ImagePath = savePath, size = (400, 400))
            PPTGenerator.replace_image(Slide = Slide, ImagePath = savePath, ShapeName = "图片")

        elif imageCount == 3:
            prompt += PageSummary["标题"]
            for i in range(3):
                imagePrompt = LLMInterface.GetResponse_String(prompt)
                imageUrl = ImageGenerator.generate_image(Prompt = imagePrompt, Size = (640, 360))
                savePath = ImageGenerator.download_image(imageUrl)
                ImageGenerator.resize_image(ImagePath = savePath, Size = (352, 198))
                PPTGenerator.replace_image(Slide = Slide, ImagePath = savePath, ShapeName = f"图片{i}")


# PPT生成器类
class PPTGenerator(SlideOperator):

    # 根据大纲生成PPT主流程
    @staticmethod
    def main_process(RequestData, UserName):
        try:
            userContent = RequestData.get("content", "")
            UUIDList = RequestData.get("materialFiles", None)
            PPTCatalog = RequestData.get("catalog", None)

            if UUIDList is None:
                raise Exception("必须传入素材文件!")
            if PPTCatalog is None:
                raise Exception("必须传入PPT大纲!")

            # 根据大纲生成内容
            pptContent = PPTGenerator.generate_content(UUIDList = UUIDList,
                                                       UserName = UserName,
                                                       Catalog = PPTCatalog,
                                                       UserContent = userContent)
            templatePath = os.path.join(GLOBAL_ResourcesSavePath, "模版.pptx")
            # 根据内容生成ppt文件
            pptObj = PPTGenerator.generate_ppt(PPTContent = pptContent, TemplatePath = templatePath)

            fileUUID = Tools.GetUUID()
            fileName = pptContent["主标题"] + ".pptx"
            # 解析PPT内容
            textContent = "主标题:" + pptContent["主标题"] + "\n" + "副标题:" + pptContent["副标题"]
            for content in pptContent["内容"]:
                chapterTitle = content["章节标题"]
                textContent += f"章节标题：{chapterTitle}\n"
                for chapterContent in content["章节内容"]:
                    pageTitle = chapterContent["页标题"]
                    textContent += f"页标题：{pageTitle}\n"
                    for pageContent in chapterContent["页内容"]:
                        sectionTitle = pageContent["节标题"]
                        sectionContent = pageContent["节内容"]
                        textContent += f"节标题{pageTitle}节内容{sectionContent}\n"

            # 获取文件数据并保存
            summaryText = LLMInterface.FileSummary(textContent)
            saveTime = Tools.GetSaveTime()
            pptObj.save(os.path.join(fileSavePath, UserName, fileUUID + ".pptx"))
            FileProcess.SaveTxt(UUID = fileUUID, Content = textContent, UserName = UserName)
            FileProcess.SaveFileInfo(FileName = fileName,
                                     Description = summaryText,
                                     SaveTime = saveTime,
                                     UserName = UserName,
                                     UUID = fileUUID)

        except Exception as e:
            curTime = Tools.GetTime()
            logging.error(f"[{curTime}]Module:[PPTMain]" + str(e))
            raise e


    # PPT生成主流程，传入完整的PPT格式json，返回PPT文件对象
    @staticmethod
    def generate_ppt(PPTContent, TemplatePath):
        try:
            prs = Presentation(TemplatePath)
            contentList = PPTContent["内容"]
            # 处理首页
            for name in ["主标题", "副标题", "汇报人"]:
                frontSlide = prs.slides[0]  # 首页
                if name == "汇报人":
                    PPTGenerator.replace_shape_text(Slide = frontSlide,
                                                    ShapeName = name,
                                                    Content = f"汇报人：{PPTContent[name]}")
                else:
                    PPTGenerator.replace_shape_text(Slide = frontSlide,
                                                    ShapeName = name,
                                                    Content = PPTContent[name])

            # 处理目录
            catelogIndex = 0  # 目录索引
            chapterCount = len(contentList)  # 章节数量
            # 遍历所有章节
            for content in contentList:
                catelogIndex += 1
                catelogSlide = prs.slides[1]
                PPTGenerator.replace_shape_text(Slide = catelogSlide,
                                                ShapeName = f"目录内容{catelogIndex}",
                                                Content = content["章节标题"])
            # 删除多余的目录信息
            for i in range(chapterCount + 1, 7):
                for shape in [i for i in catelogSlide.shapes]:
                    if shape.name == f"目录内容{i}":
                        PPTGenerator.replace_text(shape, "")
                    elif shape.name == f"目录编号{i}":
                        PPTGenerator.replace_text(shape, "")
            ############################################################
            # 处理当前章节首页
            # 遍历文本内容中的所有章节
            for curChapter in range(chapterCount):
                # 计算ppt页的初始下标
                initIndex = 2 + curChapter * 4
                # 设置当前章节的标题
                chapterInfo = contentList[curChapter]  # 内容列表元素
                PPTGenerator.replace_shape_text(Slide = prs.slides[initIndex],
                                                ShapeName = f"章节标题",
                                                Content = chapterInfo["章节标题"])
                # 处理章节中3页内容页
                pageIndex = -1
                for index in range(initIndex + 1, initIndex + 4):
                    pageIndex += 1
                    # print(len(contentList))
                    curSlide = prs.slides[index]  # 当前页的对象
                    pageList = chapterInfo["章节内容"]  # 章节内容的页信息列表
                    curPageInfo = pageList[pageIndex]  # 当前页的信息
                    curPageSummary = {"标题": "页标题：" + curPageInfo["页标题"] + "\n",
                                      "内容": []}
                    PPTGenerator.replace_shape_text(Slide = curSlide,
                                                    ShapeName = "页标题",
                                                    Content = curPageInfo["页标题"])
                    # 处理该页中的3小节内容
                    section = 0
                    for pageContent in curPageInfo["页内容"]:
                        section += 1
                        title = pageContent["节标题"]
                        text = pageContent["节内容"]
                        curPageSummary["内容"].append(f"节标题: {title}，节内容: {text}\n")
                        PPTGenerator.replace_shape_text(Slide = curSlide,
                                                        ShapeName = f"节标题{section}",
                                                        Content = title)
                        PPTGenerator.replace_shape_text(Slide = curSlide,
                                                        ShapeName = f"节内容{section}",
                                                        Content = text)
                    # 为该页PPT插入图片
                    PPTGenerator.insert_image(Slide = curSlide, PageSummary = curPageSummary)

            # 删除多余的内容页
            startIndex = 2 + chapterCount * 4
            for i in range(25, startIndex - 1, -1):
                PPTGenerator.delete_slide(prs, i)

            curTime = Tools.GetTime()
            logging.info(f"[{curTime}]PPT file generate successed.")
            return prs

        except Exception as e:
            curTime = Tools.GetTime()
            logging.error(f"[{curTime}]Module:[GeneratePPT]" + str(e))
            raise e


    # 根据用户输入内容与文件列表生成PPT大纲的json结构
    @staticmethod
    def generate_catalog(UUIDList, UserName, UserContent = ""):
        try:
            promptPath = os.path.join(GLOBAL_ResourcesSavePath, "PPT大纲.txt")
            prompt = FileProcess.ReadTxt(promptPath)

            material = UserContent + "\n"
            for UUID in UUIDList:
                material += FileProcess.ReadTxt(os.path.join(fileSavePath, UserName, UUID + ".txt")) + "\n"
            # 限制素材字数
            if len(material) > 2500:
                material = material[:2500]

            response = LLMInterface.GetResponse_String(prompt + material)[8:-4]
            catalog = json.loads(response)
            catalog["汇报人"] = f"{UserName}"
            curTime = Tools.GetTime()
            logging.info(f"[{curTime}]PPT catalog generation finished")
            return catalog

        except Exception as e:
            curTime = Tools.GetTime()
            logging.error(f"[{curTime}]Module:[GenPPTCatalog]" + str(e))
            raise e


    # 生成PPT内容json
    @staticmethod
    def generate_content(UUIDList, UserName, Catalog, UserContent = ""):
        try:
            promptPath = os.path.join(GLOBAL_ResourcesSavePath, "PPT内容.txt")
            prompt = FileProcess.ReadTxt(promptPath)

            material = UserContent + "\n"
            for UUID in UUIDList:
                material += FileProcess.ReadTxt(os.path.join(fileSavePath, UserName, UUID + ".txt")) + "\n"
            # 限制素材字数
            if len(material) > 2500:
                material = material[:2500]

            response = Catalog
            contentList = []
            for chapter in response["内容"]:
                sectionContent = ""
                cnt = 0
                chapterTitle = chapter["章节标题"]
                for pageTitle in chapter["章节内容"]:
                    cnt += 1
                    sectionContent += str(cnt) + "." + pageTitle + "\n"
                chapterContent = f"章节标题：{chapterTitle}\n页标题:\n{sectionContent}\n"
                sep = "[用户输入内容]" + "\n" + material
                finalPrompt = prompt + chapterContent + sep
                resp = LLMInterface.GetResponse_String(finalPrompt)[8:-4]
                contentList.append(json.loads(resp))

            response["内容"] = contentList
            curTime = Tools.GetTime()
            logging.info(f"[{curTime}]PPT content generation finished")
            return response

        except Exception as e:
            curTime = Tools.GetTime()
            logging.error(f"[{curTime}]Module:[GenPPTContent]" + str(e))
            raise e
